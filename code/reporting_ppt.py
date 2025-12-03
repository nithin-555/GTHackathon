from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
import pandas as pd
import numpy as np
import seaborn as sns

class PPTReporter:
    """
    Generates professional PowerPoint reports with visualizations.
    """
    
    def __init__(self):
        self.colors = {
            'primary': RGBColor(0, 102, 204),
            'secondary': RGBColor(255, 153, 0),
            'success': RGBColor(52, 152, 219),
            'text': RGBColor(51, 51, 51),
            'light_bg': RGBColor(240, 240, 240)
        }
    
    def _generate_narrative_summary(self, summary_stats, df):
        """Generate natural language narrative about the dataset"""
        overview = summary_stats.get('dataset_overview', {})
        col_info = summary_stats.get('column_information', {})
        
        total_rows = overview.get('total_rows', 0)
        total_cols = overview.get('total_columns', 0)
        
        # Identify column types
        numeric_cols = [col for col, info in col_info.items() if 'statistics' in info]
        categorical_cols = [col for col, info in col_info.items() if 'statistics' not in info]
        
        # Build narrative
        narrative = f"This comprehensive dataset encompasses {total_rows:,} observations across {total_cols} dimensions. "
        
        # Describe what the data represents
        if df is not None and len(numeric_cols) > 0:
            sample_metrics = numeric_cols[:3]
            narrative += f"Key metrics under analysis include {', '.join(sample_metrics)}. "
        
        # Data quality
        null_percentages = [info.get('null_percentage', 0) for info in col_info.values()]
        avg_completeness = 100 - (sum(null_percentages) / len(null_percentages) if null_percentages else 0)
        
        if avg_completeness > 95:
            narrative += f"The dataset demonstrates exceptional quality with {avg_completeness:.1f}% completeness. "
        elif avg_completeness > 80:
            narrative += f"Data quality is robust with {avg_completeness:.1f}% completeness across all attributes. "
        
        narrative += "Through rigorous analysis, we have identified key patterns, trends, and actionable insights that can drive strategic decision-making."
        
        return narrative
    
    def generate_report(self, summary_stats, insights, output_path, df=None):
        """
        Creates a professional multi-slide PPTX file with visualizations.
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # 1. Title Slide
            self._create_title_slide(prs)
            
            # 2. Executive Summary with narrative
            self._create_executive_summary(prs, summary_stats, df)
            
            # 3. Data visualizations
            if df is not None:
                self._create_visualization_slides(prs, df, summary_stats)
            
            # 4. AI Insights (split into sections)
            self._create_insights_slides(prs, insights)
            
            # 5. Recommendations & Next Steps
            self._create_recommendations_slide(prs, insights)
            
            os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
            prs.save(output_path)
            print(f"✓ PowerPoint Report generated at {output_path}")
            return True
        except Exception as e:
            print(f"Error generating PPT: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def _create_title_slide(self, prs):
        """Create professional title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        shape.line.fill.background()
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Executive Data Analysis"
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Strategic Insights & Recommendations"
        
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors['text']
        p.alignment = PP_ALIGN.CENTER
        
        # Decorative line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3), Inches(4.6),
            Inches(4), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['primary']
        line.line.fill.background()
    
    def _create_executive_summary(self, prs, summary_stats, df):
        """Create executive summary with natural language narrative"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Executive Overview"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Natural language narrative
        narrative = self._generate_narrative_summary(summary_stats, df)
        
        desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(2.2))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_frame.text = narrative
        
        p = desc_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors['text']
        p.line_spacing = 1.5
        p.alignment = PP_ALIGN.JUSTIFY
        
        # Key metrics cards
        overview = summary_stats.get('dataset_overview', {})
        col_info = summary_stats.get('column_information', {})
        
        numeric_count = sum(1 for v in col_info.values() if 'statistics' in v)
        
        metrics = [
            ("Total Records", f"{overview.get('total_rows', 'N/A'):,}"),
            ("Data Points", f"{overview.get('total_columns', 'N/A')} fields"),
            ("Numeric Metrics", f"{numeric_count}")
        ]
        
        box_width = Inches(2.5)
        box_height = Inches(1.4)
        start_left = Inches(1)
        top = Inches(4.5)
        spacing = Inches(0.5)
        
        for idx, (label, value) in enumerate(metrics):
            left = start_left + idx * (box_width + spacing)
            
            # Background card
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors['primary']
            shape.line.color.rgb = self.colors['primary']
            shape.shadow.inherit = False
            
            # Value
            value_box = slide.shapes.add_textbox(left, top + Inches(0.2), box_width, Inches(0.6))
            value_frame = value_box.text_frame
            value_frame.text = value
            p = value_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Label
            label_box = slide.shapes.add_textbox(left, top + Inches(0.85), box_width, Inches(0.4))
            label_frame = label_box.text_frame
            label_frame.text = label
            p = label_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
    
    def _create_visualization_slides(self, prs, df, summary_stats):
        """Create slides with professional matplotlib visualizations"""
        
        # Slide 1: Numeric Distribution Analysis
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            viz_path = self._create_numeric_charts(df, numeric_cols)
            if viz_path and os.path.exists(viz_path):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
                title_frame = title_box.text_frame
                title_frame.text = "Quantitative Analysis"
                p = title_frame.paragraphs[0]
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = self.colors['primary']
                
                slide.shapes.add_picture(viz_path, Inches(0.5), Inches(1.3), width=Inches(9))
                os.remove(viz_path)
        
        # Slide 2: Categorical Distribution
        categorical_cols = df.select_dtypes(include=['object']).columns
        if len(categorical_cols) > 0:
            viz_path = self._create_categorical_charts(df, categorical_cols)
            if viz_path and os.path.exists(viz_path):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
                title_frame = title_box.text_frame
                title_frame.text = "Categorical Distribution"
                p = title_frame.paragraphs[0]
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = self.colors['primary']
                
                slide.shapes.add_picture(viz_path, Inches(0.5), Inches(1.3), width=Inches(9))
                os.remove(viz_path)
        
        # Slide 3: Correlation Analysis (if applicable)
        if len(numeric_cols) >= 2:
            viz_path = self._create_correlation_heatmap(df, numeric_cols)
            if viz_path and os.path.exists(viz_path):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
                title_frame = title_box.text_frame
                title_frame.text = "Correlation Matrix"
                p = title_frame.paragraphs[0]
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = self.colors['primary']
                
                slide.shapes.add_picture(viz_path, Inches(1.5), Inches(1.3), width=Inches(7))
                os.remove(viz_path)
    
    def _create_numeric_charts(self, df, numeric_cols):
        """Create professional distribution charts for numeric columns"""
        try:
            cols_to_plot = list(numeric_cols[:6])
            n_cols = min(3, len(cols_to_plot))
            n_rows = (len(cols_to_plot) + n_cols - 1) // n_cols
            
            fig, axes = plt.subplots(n_rows, n_cols, figsize=(13, 3.8*n_rows))
            fig.patch.set_facecolor('white')
            
            if n_rows == 1 and n_cols == 1:
                axes = np.array([axes])
            elif n_rows == 1:
                axes = np.array([axes])
            else:
                axes = axes.flatten()
            
            colors = ['#3498db', '#e74c3c', '#2ecc71', '#f39c12', '#9b59b6', '#1abc9c']
            
            for idx, col in enumerate(cols_to_plot):
                ax = axes[idx] if len(cols_to_plot) > 1 else axes[0]
                
                data = df[col].dropna()
                
                # Histogram
                n, bins, patches = ax.hist(data, bins=25, color=colors[idx % len(colors)], 
                                           alpha=0.7, edgecolor='black', density=True)
                
                # KDE overlay
                try:
                    data.plot(kind='kde', ax=ax, color='darkred', linewidth=2.5, alpha=0.8)
                except:
                    pass
                
                ax.set_title(f'{col}', fontsize=12, fontweight='bold', pad=12)
                ax.set_xlabel('')
                ax.set_ylabel('Density', fontsize=10)
                ax.grid(axis='y', alpha=0.3, linestyle='--', linewidth=0.7)
                
                # Statistics
                mean_val = data.mean()
                median_val = data.median()
                ax.axvline(mean_val, color='red', linestyle='--', linewidth=2, alpha=0.8, label=f'Mean: {mean_val:.2f}')
                ax.axvline(median_val, color='green', linestyle='--', linewidth=2, alpha=0.8, label=f'Median: {median_val:.2f}')
                ax.legend(fontsize=8, loc='upper right')
            
            # Hide extra subplots
            for idx in range(len(cols_to_plot), len(axes)):
                fig.delaxes(axes[idx])
            
            plt.tight_layout()
            temp_path = 'temp_ppt_numeric.png'
            plt.savefig(temp_path, dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            return temp_path
        except Exception as e:
            print(f"Error creating numeric charts: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _create_categorical_charts(self, df, categorical_cols):
        """Create professional bar charts for categorical columns"""
        try:
            cols_to_plot = [col for col in categorical_cols if 2 <= df[col].nunique() <= 15][:3]
            if not cols_to_plot:
                return None
            
            fig, axes = plt.subplots(1, len(cols_to_plot), figsize=(5*len(cols_to_plot), 4.5))
            fig.patch.set_facecolor('white')
            
            if len(cols_to_plot) == 1:
                axes = [axes]
            
            colors_palette = ['#3498db', '#e74c3c', '#2ecc71']
            
            for idx, col in enumerate(cols_to_plot):
                value_counts = df[col].value_counts().head(8)
                
                bars = axes[idx].barh(range(len(value_counts)), value_counts.values, 
                                      color=colors_palette[idx % len(colors_palette)],
                                      edgecolor='black', alpha=0.85, height=0.7)
                
                axes[idx].set_yticks(range(len(value_counts)))
                axes[idx].set_yticklabels(value_counts.index, fontsize=10)
                axes[idx].set_title(f'{col}', fontsize=12, fontweight='bold', pad=12)
                axes[idx].set_xlabel('Frequency', fontsize=10)
                axes[idx].grid(axis='x', alpha=0.3, linestyle='--', linewidth=0.7)
                
                # Value labels
                for i, v in enumerate(value_counts.values):
                    axes[idx].text(v + max(value_counts.values)*0.02, i, f'{v:,}', 
                                 va='center', fontsize=9, fontweight='bold')
            
            plt.tight_layout()
            temp_path = 'temp_ppt_categorical.png'
            plt.savefig(temp_path, dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            return temp_path
        except Exception as e:
            print(f"Error creating categorical charts: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _create_correlation_heatmap(self, df, numeric_cols):
        """Create professional correlation heatmap"""
        try:
            cols = list(numeric_cols[:10])
            corr_matrix = df[cols].corr()
            
            fig, ax = plt.subplots(figsize=(9, 7))
            fig.patch.set_facecolor('white')
            
            sns.heatmap(corr_matrix, annot=True, fmt='.2f', cmap='RdYlBu_r', 
                       center=0, square=True, linewidths=1.5, 
                       cbar_kws={"shrink": 0.8}, ax=ax, vmin=-1, vmax=1,
                       annot_kws={'size': 9, 'weight': 'bold'})
            
            ax.set_title('Correlation Analysis', fontsize=16, fontweight='bold', pad=20)
            plt.xticks(rotation=45, ha='right', fontsize=10)
            plt.yticks(rotation=0, fontsize=10)
            
            plt.tight_layout()
            temp_path = 'temp_ppt_correlation.png'
            plt.savefig(temp_path, dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            return temp_path
        except Exception as e:
            print(f"Error creating correlation heatmap: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _create_insights_slides(self, prs, insights):
        """Create AI insights slides with professional formatting"""
        if not insights or insights.startswith("AI Analysis skipped"):
            return
        
        # Split insights into major sections
        sections = []
        current_section = []
        
        for line in insights.split('\n'):
            if line.strip() and (line.strip()[0].isdigit() or line.strip().startswith('#')) and any(keyword in line for keyword in ['Key Statistical', 'Notable Patterns', 'Anomalies', 'Business Implications', 'Recommended']):
                if current_section:
                    sections.append('\n'.join(current_section))
                current_section = [line]
            else:
                current_section.append(line)
        
        if current_section:
            sections.append('\n'.join(current_section))
        
        # Create slides for major sections (excluding Recommended Actions)
        for section in sections:
            if 'Recommended Actions' in section:
                continue
                
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Extract title
            first_line = section.split('\n')[0]
            title_text = first_line.replace('#', '').replace('*', '').strip()
            if title_text[0].isdigit():
                title_text = '. '.join(title_text.split('.')[1:]).strip()
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = title_text
            p = title_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            # Content
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.7))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            lines = section.split('\n')[1:]
            for line in lines:
                if not line.strip():
                    continue
                
                line = line.replace('**', '').replace('##', '').strip()
                
                if text_frame.paragraphs[0].text:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                
                if line.startswith('-') or line.startswith('•'):
                    p.text = f"• {line[1:].strip()}"
                    p.level = 0
                else:
                    p.text = line
                
                # Formatting
                if ':' in line and len(line) < 80 and not line.startswith('•'):
                    p.font.bold = True
                    p.font.size = Pt(15)
                else:
                    p.font.size = Pt(14)
                
                p.font.color.rgb = self.colors['text']
                p.space_after = Pt(12)
                p.line_spacing = 1.3
    
    def _create_recommendations_slide(self, prs, insights):
        """Create professional recommendations slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background accent
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Strategic Recommendations"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        if "Recommended Actions" in insights:
            lines = insights.split('\n')
            in_recommendations = False
            
            for line in lines:
                if "Recommended Actions" in line:
                    in_recommendations = True
                    continue
                
                if in_recommendations and line.strip():
                    if line.strip()[0].isdigit() and '.' in line[:3] and "Recommended" not in line:
                        break
                    
                    if line.strip().startswith('-') or line.strip().startswith('•') or line.strip()[0].isdigit():
                        p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
                        clean_line = line.replace('**', '').replace('-', '').replace('•', '').strip()
                        if clean_line[0].isdigit():
                            clean_line = '. '.join(clean_line.split('.')[1:]).strip()
                        p.text = f"• {clean_line}"
                        p.font.size = Pt(16)
                        p.font.color.rgb = self.colors['text']
                        p.space_after = Pt(16)
                        p.line_spacing = 1.4
        else:
            recommendations = [
                "Leverage key statistical findings to inform strategic decisions",
                "Focus on high-impact patterns identified in the analysis",
                "Address any data quality issues to improve future insights",
                "Implement a regular review cycle for continuous improvement",
                "Share insights across teams to align on data-driven strategy"
            ]
            
            for rec in recommendations:
                p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
                p.text = f"• {rec}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.colors['text']
                p.space_after = Pt(16)
                p.line_spacing = 1.4